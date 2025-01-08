import os
import pandas as pd # type: ignore
import matplotlib.pyplot as plt # type: ignore
import re
from pptx import Presentation # type: ignore
from pptx.util import Inches # type: ignore

def process_csv_file(file_path, output_directory):
    # 读取 CSV 文件
    df = pd.read_csv(file_path)
    
    # 检查是否包含 pos_y 列
    if 'posY' not in df.columns:
        print(f'Error: {file_path} does not contain pos_y column.')
        return
    
    # 提取 pos_y 列数据
    pos_y_data = df['posY']
    
    # 打印当前正在处理的文件名
    print(f'Processing {file_path}')
    
    # 绘制图表
    plt.figure()
    plt.plot(pos_y_data)
    plt.title('Position Y over Frames')
    plt.xlabel('Index')
    plt.ylabel('Position Y')
    
    # 构建输出文件名
    file_name = os.path.basename(file_path)
    output_file = os.path.splitext(file_name)[0] + '.png'
    output_path = os.path.join(output_directory, output_file)
    
    # 保存图表为文件
    plt.savefig(output_path)
    plt.close()  # 关闭当前图表以便进行下一个图的绘制
    
    print(f'Saved plot to {output_path}')


def extract_id(file_name):
    match = re.search(r'walk_test_(\d+)_',file_name)
    if match:
            print(f'id=',match.group(1))
            #print(match.group(1))
            return match.group(1)
    return None



def process_all_csv_files(directory, output_directory):
    # 遍历目录中的所有文件
    current_folder =os.getcwd()
    directory = current_folder + directory 
    output_directory = current_folder + output_directory

    for filename in os.listdir(directory):
        if filename.endswith('.csv'):
            file_path = os.path.join(directory, filename)
            process_csv_file(file_path, output_directory)
            extract_id(filename)


def make_pptx(png_foler):
     
    current_folder =os.getcwd()
    png_foler = current_folder + png_foler

    prs = Presentation()

    files = os.listdir(png_foler)
    sorted_files = sorted(files)
     
    #for png_file in os.listdir(png_foler):
    for png_file in sorted_files:
            # 提取標題（去除 .png 後綴）
        title = png_file.rsplit('.', 1)[0]

        # 添加一頁投影片
        slide_layout = prs.slide_layouts[5]  # 空白投影片佈局
        slide = prs.slides.add_slide(slide_layout)

        # 添加標題
        title_shape = slide.shapes.title
        title_shape.text = title

        # 在投影片中加入圖片
        left = Inches(1)
        top = Inches(1.5)
        
        slide.shapes.add_picture(png_foler + "/"+png_file, left, top, height=Inches(5))
        #slide.shapes.add_picture(png_file, left, top, height=Inches(5))

        # 保存演示文稿
    
    prs.save(png_foler + 'presentation.pptx')

# 指定包含 CSV 文件的目录路径     rawdata_folder
input_directory = '/rawdata_folder'
# 指定保存 PNG 文件的目录路径
output_directory = '/png_folder'

# 确保输出目录存在
#os.makedirs(output_directory, exist_ok=True)

# 处理所有 CSV 文件
current_folder =os.getcwd()
process_all_csv_files(input_directory, output_directory)
make_pptx(output_directory)